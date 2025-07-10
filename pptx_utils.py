import pandas as pd
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io
import base64
from PIL import Image
import tempfile
import os
import warnings
warnings.filterwarnings('ignore')

def figure_to_image(fig, width=800, height=600):
    """
    Convert a Plotly figure to a PIL Image
    
    Args:
        fig: plotly.graph_objects.Figure
        width: image width in pixels
        height: image height in pixels
        
    Returns:
        PIL.Image: The figure as an image
    """
    # Convert figure to image bytes
    img_bytes = fig.to_image(format="png", width=width, height=height)
    
    # Convert bytes to PIL Image
    img = Image.open(io.BytesIO(img_bytes))
    return img

def create_title_slide(prs, title="ExcelInsight Analysis"):
    """
    Create a title slide for the presentation
    
    Args:
        prs: Presentation object
        title: Title for the presentation
        
    Returns:
        slide: The created slide
    """
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Style the title
    title_frame = title_shape.text_frame
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(18, 85, 181)  # #1255b5
    
    # Add subtitle
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = "Data Analysis & Visualization"
    
    return slide

def create_chart_slide(prs, fig, chart_info, slide_number):
    """
    Create a slide with a chart
    
    Args:
        prs: Presentation object
        fig: plotly.graph_objects.Figure
        chart_info: dictionary with chart information
        slide_number: slide number for title
        
    Returns:
        slide: The created slide
    """
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = f"Chart {slide_number}: {chart_info['type'].title()} Chart"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(18, 85, 181)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add subtitle with column information
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = f"Columns: {', '.join(chart_info['columns'])}"
    subtitle_frame.paragraphs[0].font.size = Pt(14)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Convert figure to image
    img = figure_to_image(fig, width=800, height=500)
    
    # Save image to temporary file
    with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
        img_path = tmp_file.name
        img.save(img_path)
    
    # Add image to slide
    try:
        slide.shapes.add_picture(img_path, Inches(1), Inches(2), Inches(8), Inches(5))
    finally:
        # Clean up temporary file
        os.unlink(img_path)
    
    return slide

def create_data_profiling_slide(prs, df):
    """
    Create a slide with data profiling information
    
    Args:
        prs: Presentation object
        df: pandas DataFrame
        
    Returns:
        slide: The created slide
    """
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Data Profiling Summary"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(18, 85, 181)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Create profiling information
    profiling_info = []
    profiling_info.append(f"📊 Dataset Overview")
    profiling_info.append(f"• Total Rows: {len(df):,}")
    profiling_info.append(f"• Total Columns: {len(df.columns)}")
    profiling_info.append(f"• Memory Usage: {df.memory_usage(deep=True).sum() / 1024:.1f} KB")
    profiling_info.append("")
    
    # Column information
    profiling_info.append("📋 Column Analysis")
    for col in df.columns:
        col_type = str(df[col].dtype)
        null_count = df[col].isnull().sum()
        null_pct = (null_count / len(df)) * 100
        
        profiling_info.append(f"• {col}: {col_type}")
        profiling_info.append(f"  - Null values: {null_count} ({null_pct:.1f}%)")
        
        if pd.api.types.is_numeric_dtype(df[col]):
            profiling_info.append(f"  - Range: {df[col].min():.2f} to {df[col].max():.2f}")
        elif pd.api.types.is_datetime64_any_dtype(df[col]):
            profiling_info.append(f"  - Date range: {df[col].min()} to {df[col].max()}")
        else:
            profiling_info.append(f"  - Unique values: {df[col].nunique()}")
    
    # Add profiling text to slide
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(6))
    text_frame = text_box.text_frame
    text_frame.text = "\n".join(profiling_info)
    
    # Style the text
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(12)
        paragraph.font.color.rgb = RGBColor(50, 50, 50)
    
    return slide

def create_powerpoint_deck(df, chart_candidates, output_path, include_profiling=False):
    """
    Create a PowerPoint presentation with charts and optional data profiling
    
    Args:
        df: pandas DataFrame
        chart_candidates: list of chart candidate dictionaries
        output_path: path to save the PowerPoint file
        include_profiling: whether to include data profiling slide
    """
    # Create presentation
    prs = Presentation()
    
    # Set slide dimensions (16:9 aspect ratio)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Create title slide
    create_title_slide(prs)
    
    # Add data profiling slide if requested
    if include_profiling:
        create_data_profiling_slide(prs, df)
    
    # Create chart slides
    for i, chart_info in enumerate(chart_candidates):
        try:
            # Import here to avoid circular imports
            from chart_utils import create_chart
            
            # Create the chart
            fig = create_chart(df, chart_info)
            
            # Create slide with chart
            create_chart_slide(prs, fig, chart_info, i + 1)
            
        except Exception as e:
            print(f"Warning: Could not create chart {i+1}: {str(e)}")
            continue
    
    # Save the presentation
    prs.save(output_path) 